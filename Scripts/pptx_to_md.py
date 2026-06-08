#!/usr/bin/env python3
"""Convert PowerPoint (.pptx) to Markdown."""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from Tools.drawio_writer import from_json, write_drawio
from Tools.markdown_utils import build_frontmatter, ensure_output_dir


def convert(input_path: Path, output_dir: Path) -> Path:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    input_path = Path(input_path)
    basename = input_path.stem
    out_dir = ensure_output_dir(output_dir, basename)
    assets_dir = out_dir / f"{basename}_assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    charts_dir = out_dir / "charts"
    charts_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(input_path))
    parts: list[str] = []
    img_counter = 0
    chart_counter = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()
        heading = f"Slide {slide_num}"
        if title:
            heading += f": {title}"
        parts.append(f"## {heading}")
        parts.append("")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and shape != slide.shapes.title:
                        parts.append(text)
                        parts.append("")

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_counter += 1
                try:
                    img = shape.image
                    ext = img.ext or "png"
                    img_name = f"slide{slide_num}_img{img_counter}.{ext}"
                    img_path = assets_dir / img_name
                    img_path.write_bytes(img.blob)
                    parts.append(f"![Slide {slide_num} image]({assets_dir.name}/{img_name})")
                    parts.append("")
                except Exception:
                    parts.append(f"<!-- Image on slide {slide_num} could not be extracted -->")
                    parts.append("")

            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart_counter += 1
                chart_name = f"slide{slide_num}_chart{chart_counter}"
                chart_json = {
                    "name": chart_name,
                    "nodes": [
                        {"key": "c", "label": f"Chart on slide {slide_num}\n(Edit in draw.io)"},
                    ],
                    "edges": [],
                }
                write_drawio(from_json(chart_json), charts_dir / f"{chart_name}.drawio")
                parts.append(f"[Edit chart in draw.io](charts/{chart_name}.drawio)")
                parts.append("")

    md_path = out_dir / f"{basename}.md"
    frontmatter = build_frontmatter(
        title=basename,
        source=str(input_path),
        converter="pptx_to_md",
    )
    md_path.write_text(frontmatter + "\n".join(parts), encoding="utf-8")
    return md_path


def main() -> int:
    parser = argparse.ArgumentParser(description="PowerPoint to Markdown")
    parser.add_argument("--input", "-i", required=True)
    parser.add_argument("--output", "-o", default="output")
    args = parser.parse_args()

    try:
        out = convert(Path(args.input), Path(args.output))
        print(out)
        return 0
    except FileNotFoundError as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 2
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 2


if __name__ == "__main__":
    raise SystemExit(main())
